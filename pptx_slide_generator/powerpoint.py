from typing import Dict, Tuple, List, Optional

from loguru import logger

from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from pptx_slide_generator.models import ExcelData, RoomData, SlideData

logger.add("logs/file_{time}.log", level="WARNING", format="{message}")


def _get_relevant_shapes_by_name(
        slide: Slide,
        exclude_name: Optional[str] = None,
        include_prefix: Optional[str] = None) -> List[BaseShape]:
    """Fetches all shapes for given slide which are not ignored.

    """

    shapes = slide.shapes

    if exclude_name:
        shapes = [shape for shape in shapes
                  if shape.name != exclude_name]

    if include_prefix:
        shapes = [shape for shape in shapes
                  if shape.name.startswith(include_prefix)]

    return shapes


def _get_sanitized_room_from_shape_name(shape: BaseShape, prefix: str):
    """Get the lookup key between pptx shape and excel row while sanitizing.

    """

    return shape.name.replace(prefix, "").replace("_", "/")


def get_named_master_layouts(pptx: Presentation) -> Dict:
    """Get master layouts with names as unique identifiers.

    """

    slide_layouts = {}

    for slide_master in pptx.slide_masters:
        for slide_layout in slide_master.slide_layouts:
            assert slide_layout.name not in slide_layouts
            slide_layouts[slide_layout.name] = slide_layout

    return slide_layouts


def rename_shapes(pptx: Presentation,
                  slide_idx: int,
                  names: Tuple[str],
                  exclude: str):
    """Renames shapes for a given slide index while ignoring shapes with a
    specific name.

    """

    slide = pptx.slides[slide_idx]
    shapes = _get_relevant_shapes_by_name(slide, exclude_name=exclude)
    assert len(shapes) == len(names)

    shape_names = zip(shapes, names)
    for shape, name in shape_names:
        shape.name = name


def _get_shape_coordinates(shape: BaseShape) -> Tuple[float, float]:
    """Get top and left coordinates of shape."""

    return shape.top, shape.left


def _mirror_shape_names(source: Slide, target: Slide):
    """Mirrors all shape names in source slide with target slice. Leverages
    shape coordinates to create mapping between source and target shape because
    order can't be assumed to be equal.

    """

    keyed_source_names = {_get_shape_coordinates(shape): shape.name
                          for shape in source.shapes}

    for shape in target.shapes:
        target_key = _get_shape_coordinates(shape)
        shape.name = keyed_source_names[target_key]


def _get_room_data(excel_data: ExcelData, room_name: str) -> Optional[RoomData]:
    """Retrieve slide data for given a room while performing validations.

    """

    try:
        data = excel_data.get(room_name)
    except KeyError:
        logger.warning(f"Room '{room_name}' from SVG/Pptx does not have any data in excel file.")
        return

    return data


def _generate_room_slides(pptx: Presentation,
                          layouts: Dict,
                          room_data: RoomData,
                          shape: BaseShape,
                          exclude: str):
    """Generates all slides for a single room.

    """

    for idx, slide_data in enumerate(room_data.slides):

        if not slide_data.relevant:
            logger.info(f"Room '{room_data.name}' from SVG/Pptx is skipped "
                        f"because marked as not relevant in excel file.")
            continue

        if not slide_data.layout:
            logger.warning(f"Room '{room_data.name}' from SVG/Pptx is skipped "
                           f"because no layout provided in excel file.")
            continue

        try:
            layout = layouts[slide_data.layout]
        except KeyError:
            logger.warning(f"Room '{room_data.name}' has no corresponding layout "
                           f"'{slide_data.layout}' in Powerpoint master.")
            continue

        new_slide = pptx.slides.add_slide(layout)

        # only add a single hyperlink for first slide because multiple slides can't be linked
        if idx == 0:
            shape.click_action.target_slide = new_slide

        _populate_shapes(layout=layout,
                         new_slide=new_slide,
                         data=slide_data,
                         ignore=exclude,
                         room=room_data.name)


def generate_slides(pptx: Presentation,
                    slide_idx: int,
                    excel_data: ExcelData,
                    prefix: str,
                    exclude: str):
    """Generates new slides given a powerpoint slide with relevant
    shape names and corresponding slide data.

    """

    slide = pptx.slides[slide_idx]
    shapes = _get_relevant_shapes_by_name(slide, include_prefix=prefix)
    shapes = sorted(shapes, key=lambda x: x.name)
    layouts = get_named_master_layouts(pptx)

    for shape in shapes:
        room_name = _get_sanitized_room_from_shape_name(shape, prefix)
        room_data = _get_room_data(excel_data=excel_data, room_name=room_name)

        if not room_data:
            continue

        _generate_room_slides(pptx=pptx,
                              layouts=layouts,
                              room_data=room_data,
                              shape=shape,
                              exclude=exclude)


def _populate_shapes(layout: Slide,
                     new_slide: Slide,
                     data: SlideData,
                     ignore: str,
                     room: str):
    """Populate newly created shapes with content from `SlideData`.

    """

    _mirror_shape_names(layout, new_slide)
    shapes = _get_relevant_shapes_by_name(new_slide, exclude_name=ignore)

    for shape in shapes:
        try:
            shape.text = data.values[shape.name]
        except KeyError:
            logger.warning(f"Room '{room}' with layout '{layout.name}' "
                           f"has no data for shape '{shape.name}' in excel.")
